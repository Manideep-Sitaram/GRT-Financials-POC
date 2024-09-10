import io
import streamlit as st
from model import load_documents_initially, user_input, reset_vector_database
import os
from dotenv import load_dotenv
import logging
from pptx import Presentation
from pptx.util import Pt

# Configure logging
logging.basicConfig(
    filename='app.log',  # Log file name
    level=logging.INFO,  # Log level (DEBUG, INFO, WARNING, ERROR, CRITICAL)
    format='%(asctime)s %(levelname)s %(message)s',  # Log format
    datefmt='%Y-%m-%d %H:%M:%S'  # Date format
)

# Create a logger
logger = logging.getLogger(__name__)

# Load environment variables from .env file
load_dotenv()

# Override default sqlite3 with pysqlite3 for better performance
__import__('pysqlite3')
import sys
sys.modules['sqlite3'] = sys.modules.pop('pysqlite3')

# Set environment variables for LangChain and OpenAI APIs
os.environ["LANGCHAIN_TRACING_V2"] = "true"
os.environ["OPENAI_API_KEY"] = st.secrets["OPENAI_API_KEY"]
os.environ["LANGCHAIN_API_KEY"] = st.secrets["LANGCHAIN_API_KEY"]
os.environ["LANGCHAIN_PROJECT"] = st.secrets["LANGCHAIN_PROJECT"]

# Streamlit application title
st.title("CMR Financials Application")

if "file_uploader_key" not in st.session_state:
    st.session_state["file_uploader_key"] = 0
    
if "messages" not in st.session_state:
    st.session_state.messages = []
    
    
def reset_state():
    st.session_state.messages = []
    st.session_state["file_uploader_key"] += 1
    st.experimental_rerun()
    # st.rerun()

def get_chat_conversation_ppt():
    prs = Presentation()
    for message in st.session_state.messages:
        if message["role"] == "assistant":
            content = message["content"]
            category = prev_category if prev_category else "Assistant"
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            body = slide.shapes.placeholders[1]
            title.text = category
                        # Set the content and font size
            text_frame = body.text_frame
            text_frame.clear()  # Clear any existing content
            text_frame.word_wrap = True
            
            p = text_frame.add_paragraph()
            p.text = content
            p.font.size = Pt(13)  # Set font size to 24 points
        elif message["role"] == "user":
           prev_category = message["content"]
    return prs
    

for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        print(message)
        st.markdown(message["content"], unsafe_allow_html=True)
        # if "pageNumbers" in message:
        #     st.markdown(f"Page Numbers : {message['pageNumbers']}")
        if "pageNumbersWithSourceName" in message:
            st.markdown(f"Page Numbers with Source Name : {message['pageNumbersWithSourceName']}")
            


# Sidebar menu
with st.sidebar:
    st.title("Menu")
    pdf_docs = st.file_uploader("Upload your PDF Files and Click on the Submit & Process", type="pdf", key=st.session_state["file_uploader_key"], accept_multiple_files=True)
    submit_process = st.button("Submit & Process")
    
    reset_button = st.button("Reset All")
    if reset_button:
        reset_state()


if user_query := st.chat_input("Please Enter Your Query"):
    with st.chat_message("user"):
        st.session_state.messages.append({"role": "user", "content": user_query,"userChat":True})
        st.markdown(user_query)
    with st.chat_message("assistant"):
        response_object = user_input(user_query)
        response = (response_object["response"]).replace("$","\$")
        # pageNumbers = response_object["pageNumbers"]
        pageNumbersWithSourceName = response_object["pageNumbersWithSourceName"]
        # st.session_state.messages.append({"role": "assistant", "content": response,"pageNumbers":pageNumbers})
        st.session_state.messages.append({"role": "assistant", "content": response,"pageNumbersWithSourceName":pageNumbersWithSourceName})
        st.markdown(response, unsafe_allow_html=True)
        # st.markdown(f"Page Numbers : {pageNumbers}")
        st.markdown(f"Page Numbers with Source Name : {pageNumbersWithSourceName}")

# Function to get the chat conversation as Markdown
def get_chat_conversation_markdown():
    conversation_md = ""
    for message in st.session_state.messages:
        role = message["role"]
        content = message["content"]
        conversation_md += f"**{role.capitalize()}:** {content}\n\n"
    return conversation_md
     
if submit_process:
    if pdf_docs:
        with st.spinner("Processing..."):
            initial_question_and_answers = load_documents_initially(pdf_docs)
            st.write("Submit & Process is clicked")

            prev_category = ''

            for qa_pair in initial_question_and_answers:
                question = qa_pair["question"]
                category = qa_pair["category"]
                # pageNumbers = qa_pair["pageNumbers"]
                pageNumbersWithSourceName = qa_pair["pageNumbersWithSourceName"]
                answer = qa_pair["answer"].replace("$","\$")

                logging.info(f"Question: {question}, Category: {category}")

                if category != prev_category:
                    st.write("---")  # This will add a horizontal line for better readability
                    st.subheader(f"{category}")
                st.session_state.messages.append({"role": "user", "content": category})
                # st.session_state.messages.append({"role": "assistant", "content": answer,"pageNumbers": pageNumbers})
                st.session_state.messages.append({"role": "assistant", "content": answer,"pageNumbersWithSourceName": pageNumbersWithSourceName})

                with st.chat_message("assistant"):
                    st.markdown(answer, unsafe_allow_html=True)
                    # st.markdown(f"Page Numbers : {pageNumbers}")
                    st.markdown(f"Page Numbers With Source Name : {pageNumbersWithSourceName}")

                prev_category = category

    else:
        st.warning("Please Upload The PDF")

# Download and Reset buttons
col1, col2 = st.columns(2)

with col1:
    prs = get_chat_conversation_ppt()
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    download_str = f"GRT_Financials_Chat_{st.session_state.get('session_id', 'default')}.pptx"
    
    st.download_button(
        label="Download Chat Conversation",
        data=pptx_bytes,
        file_name=download_str,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )